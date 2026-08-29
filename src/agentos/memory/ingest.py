"""Document extraction and ingestion for the AgentOS knowledge base memory tier."""

from __future__ import annotations

import io
from dataclasses import dataclass
from pathlib import Path
from typing import TYPE_CHECKING, Any

import structlog

from .types import MemorySource

if TYPE_CHECKING:
    from .store import LongTermMemoryStore

logger = structlog.get_logger(__name__)

SUPPORTED_EXTENSIONS: frozenset[str] = frozenset(
    {
        ".txt",
        ".md",
        ".markdown",
        ".rst",
        ".pdf",
        ".docx",
        ".pptx",
        ".json",
        ".csv",
        ".tsv",
        ".yaml",
        ".yml",
        ".toml",
        ".xml",
        ".html",
        ".htm",
        ".py",
        ".js",
        ".ts",
        ".jsx",
        ".tsx",
        ".css",
        ".sh",
        ".log",
    }
)

MAX_DOCUMENT_SIZE_BYTES: int = 10 * 1024 * 1024  # 10 MB limit for single document ingestion


@dataclass(frozen=True)
class IngestDocumentResult:
    path: str
    title: str
    size_bytes: int
    chunks_indexed: int
    status: str  # "indexed" | "unchanged" | "error"
    error: str | None = None

    def as_dict(self) -> dict[str, Any]:
        return {
            "path": self.path,
            "title": self.title,
            "sizeBytes": self.size_bytes,
            "chunksIndexed": self.chunks_indexed,
            "status": self.status,
            "error": self.error,
        }


def _extract_pdf_text(data: bytes | io.BytesIO | Path) -> str:
    """Extract readable text from a PDF file using pypdf or pdfplumber."""
    try:
        import pypdf

        stream: Any = (
            io.BytesIO(data)
            if isinstance(data, bytes)
            else (str(data) if isinstance(data, Path) else data)
        )
        reader = pypdf.PdfReader(stream)
        pages: list[str] = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            text = text.strip()
            if text:
                pages.append(f"[Page {i + 1}]\n{text}")
        if pages:
            return "\n\n".join(pages)
    except Exception as exc:  # noqa: BLE001
        logger.debug("ingest.pypdf_failed", error=str(exc))

    try:
        import pdfplumber

        stream = (
            io.BytesIO(data)
            if isinstance(data, bytes)
            else (str(data) if isinstance(data, Path) else data)
        )
        with pdfplumber.open(stream) as doc:
            pages = []
            for i, p in enumerate(doc.pages):
                text = p.extract_text() or ""
                text = text.strip()
                if text:
                    pages.append(f"[Page {i + 1}]\n{text}")
            return "\n\n".join(pages)
    except Exception as exc:  # noqa: BLE001
        logger.warning("ingest.pdfplumber_failed", error=str(exc))
        raise ValueError(f"Failed to extract text from PDF: {exc}") from exc


def _extract_docx_text(data: bytes | io.BytesIO | Path) -> str:
    """Extract text from a Word document (.docx)."""
    try:
        import docx

        stream: Any = (
            io.BytesIO(data)
            if isinstance(data, bytes)
            else (str(data) if isinstance(data, Path) else data)
        )
        doc = docx.Document(stream)
        paragraphs: list[str] = []
        for p in doc.paragraphs:
            text = p.text.strip()
            if text:
                paragraphs.append(text)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs.append(row_text)
        return "\n\n".join(paragraphs)
    except Exception as exc:  # noqa: BLE001
        logger.warning("ingest.docx_failed", error=str(exc))
        raise ValueError(f"Failed to extract text from DOCX: {exc}") from exc


def _extract_pptx_text(data: bytes | io.BytesIO | Path) -> str:
    """Extract text from a PowerPoint presentation (.pptx)."""
    try:
        import pptx

        stream: Any = (
            io.BytesIO(data)
            if isinstance(data, bytes)
            else (str(data) if isinstance(data, Path) else data)
        )
        prs = pptx.Presentation(stream)
        slides_text: list[str] = []
        for i, slide in enumerate(prs.slides):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_parts.append(text)
            if slide_parts:
                slides_text.append(f"[Slide {i + 1}]\n" + "\n".join(slide_parts))
        return "\n\n".join(slides_text)
    except Exception as exc:  # noqa: BLE001
        logger.warning("ingest.pptx_failed", error=str(exc))
        raise ValueError(f"Failed to extract text from PPTX: {exc}") from exc


def extract_document_text(
    source: str | Path | bytes,
    filename: str | None = None,
) -> str:
    """Extract text content from various file formats."""
    suffix = ""
    if filename:
        suffix = Path(filename).suffix.lower()
    elif isinstance(source, (str, Path)):
        suffix = Path(source).suffix.lower()

    if isinstance(source, (str, Path)):
        p = Path(source)
        if not p.is_file():
            raise FileNotFoundError(f"Document file not found: {p}")
        if suffix == ".pdf":
            return _extract_pdf_text(p)
        if suffix == ".docx":
            return _extract_docx_text(p)
        if suffix == ".pptx":
            return _extract_pptx_text(p)
        return p.read_text(encoding="utf-8", errors="replace")

    raw_bytes = source
    if suffix == ".pdf":
        return _extract_pdf_text(raw_bytes)
    if suffix == ".docx":
        return _extract_docx_text(raw_bytes)
    if suffix == ".pptx":
        return _extract_pptx_text(raw_bytes)
    return raw_bytes.decode(encoding="utf-8", errors="replace")


async def ingest_document(
    store: LongTermMemoryStore,
    source: str | Path | bytes,
    *,
    rel_path: str,
    title: str | None = None,
    mtime: float | None = None,
) -> IngestDocumentResult:
    """Ingest a single document into the memory store under MemorySource.knowledge_base."""
    doc_title = title or Path(rel_path).name
    size_bytes = 0
    if isinstance(source, (str, Path)):
        p = Path(source)
        if p.is_file():
            size_bytes = p.stat().st_size
            if mtime is None:
                mtime = p.stat().st_mtime
    elif isinstance(source, bytes):
        size_bytes = len(source)

    if size_bytes > MAX_DOCUMENT_SIZE_BYTES:
        return IngestDocumentResult(
            path=rel_path,
            title=doc_title,
            size_bytes=size_bytes,
            chunks_indexed=0,
            status="error",
            error=f"Document exceeds maximum size limit ({MAX_DOCUMENT_SIZE_BYTES} bytes)",
        )

    try:
        content = extract_document_text(source, filename=rel_path)
        content_clean = content.strip()
        if not content_clean:
            return IngestDocumentResult(
                path=rel_path,
                title=doc_title,
                size_bytes=size_bytes,
                chunks_indexed=0,
                status="error",
                error="Document contains no extractable text",
            )

        # Prefix with metadata header if not present
        if not content_clean.startswith("---"):
            header = f"---\ntitle: {doc_title}\nsource: knowledge_base\npath: {rel_path}\n---\n\n"
            indexed_content = header + content_clean
        else:
            indexed_content = content_clean

        chunks = await store.index_file(
            path=rel_path,
            content=indexed_content,
            source=MemorySource.knowledge_base,
            mtime=mtime,
        )

        status = "indexed" if chunks > 0 else "unchanged"
        return IngestDocumentResult(
            path=rel_path,
            title=doc_title,
            size_bytes=size_bytes,
            chunks_indexed=chunks,
            status=status,
        )
    except Exception as exc:  # noqa: BLE001
        logger.warning("ingest.document_failed", path=rel_path, error=str(exc))
        return IngestDocumentResult(
            path=rel_path,
            title=doc_title,
            size_bytes=size_bytes,
            chunks_indexed=0,
            status="error",
            error=str(exc),
        )


async def ingest_directory(
    store: LongTermMemoryStore,
    directory_path: str | Path,
    *,
    base_rel_prefix: str = "knowledge_base",
    recursive: bool = True,
) -> list[IngestDocumentResult]:
    """Recursively scan a directory and ingest all supported documents into the store."""
    root = Path(directory_path).resolve()
    if not root.is_dir():
        raise NotADirectoryError(f"Directory not found: {root}")

    results: list[IngestDocumentResult] = []
    pattern = "**/*" if recursive else "*"
    for item in sorted(root.glob(pattern)):
        if not item.is_file():
            continue
        if any(part.startswith(".") for part in item.relative_to(root).parts):
            continue
        suffix = item.suffix.lower()
        if suffix not in SUPPORTED_EXTENSIONS and not suffix == "":
            continue

        rel_within_dir = item.relative_to(root).as_posix()
        rel_path = f"{base_rel_prefix}/{rel_within_dir}".replace("//", "/")
        res = await ingest_document(
            store,
            item,
            rel_path=rel_path,
            title=item.name,
            mtime=item.stat().st_mtime,
        )
        results.append(res)

    return results
